import os
from typing import Dict, Any
# Import libraries inside functions or try-except blocks if optional, 
# but here they are required.
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None
    
try:
    import pptx
except ImportError:
    pptx = None

def read_content(path: str) -> Dict[str, str]:
    """
    Analyzes the file content to extract metadata for classification.
    Returns a dictionary:
    {
        "title": "Detected title or empty",
        "subtitle": "Detected subtitle or empty",
        "sample_text": "Extracted text content for keyword searching"
    }
    """
    ext = os.path.splitext(path)[1].lower()
    
    metadata = {
        "title": "",
        "subtitle": "",
        "sample_text": ""
    }
    
    try:
        if ext == '.pdf':
            metadata = _read_pdf(path)
        elif ext == '.docx':
            metadata = _read_docx(path)
        elif ext == '.xlsx':
            metadata = _read_xlsx(path)
        elif ext == '.pptx':
            metadata = _read_pptx(path)
    except Exception as e:
        # If reading fails, return empty metadata but don't crash
        # Maybe log error?
        print(f"Error reading {path}: {e}")
        pass
        
    return metadata

def _read_pdf(path: str) -> Dict[str, str]:
    if not pypdf:
        return {"title": "", "subtitle": "", "sample_text": ""}
        
    text_content = []
    title = ""
    
    with open(path, 'rb') as f:
        reader = pypdf.PdfReader(f)
        meta = reader.metadata
        if meta and meta.title:
            title = meta.title
            
        # Extract text from first few pages
        max_pages = min(len(reader.pages), 3)
        for i in range(max_pages):
            page_text = reader.pages[i].extract_text()
            if page_text:
                text_content.append(page_text)
                
    return {
        "title": title,
        "subtitle": "",
        "sample_text": "\n".join(text_content)[:2000] # Limit sample size
    }

def _read_docx(path: str) -> Dict[str, str]:
    if not docx:
        return {"title": "", "subtitle": "", "sample_text": ""}
        
    doc = docx.Document(path)
    text_content = []
    title = ""
    subtitle = ""
    
    # Try to find title style or just use first paragraphs
    for para in doc.paragraphs:
        style_name = para.style.name.lower() if para.style else ""
        text = para.text.strip()
        if not text:
            continue
            
        if 'title' in style_name and not title:
            title = text
        elif 'subtitle' in style_name and not subtitle:
            subtitle = text
        else:
            text_content.append(text)
            
        if len(text_content) > 20: # Read first ~20 paragraphs
            break
            
    return {
        "title": title,
        "subtitle": subtitle,
        "sample_text": "\n".join(text_content)[:2000]
    }

def _read_xlsx(path: str) -> Dict[str, str]:
    if not openpyxl:
        return {"title": "", "subtitle": "", "sample_text": ""}
        
    # Use context manager to ensure file is closed, preventing lock issues
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        try:
            text_content = []
            
            # Read first sheet
            if wb.sheetnames:
                ws = wb[wb.sheetnames[0]]
                # Read first 20 rows
                for i, row in enumerate(ws.iter_rows(max_row=20, values_only=True)):
                    row_text = " ".join([str(c) for c in row if c is not None])
                    if row_text.strip():
                        text_content.append(row_text)
        finally:
             wb.close()
    except Exception:
        # If openpyxl fails completely (e.g. invalid file), just return empty
        return {"title": "", "subtitle": "", "sample_text": ""}

    return {
        "title": "",
        "subtitle": "",
        "sample_text": "\n".join(text_content)[:2000]
    }

def _read_pptx(path: str) -> Dict[str, str]:
    if not pptx:
        return {"title": "", "subtitle": "", "sample_text": ""}
        
    prs = pptx.Presentation(path)
    text_content = []
    title = ""
    
    # Iterate through slides
    for i, slide in enumerate(prs.slides):
        if i > 5: break # First 5 slides
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if not text: continue
                
                # Heuristic: title of first slide is likely title
                if i == 0 and not title:
                    # Check if it's a title placeholder?
                    # shape.is_placeholder -> shape.placeholder_format.type == TITLE
                    try:
                        if shape.is_placeholder and shape.placeholder_format.type == 1: # 1 is Title
                            title = text
                    except:
                        pass
                    
                    if not title: # Fallback
                        title = text
                
                text_content.append(text)
                
    return {
        "title": title,
        "subtitle": "",
        "sample_text": "\n".join(text_content)[:2000]
    }
